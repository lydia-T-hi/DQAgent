"""
DQAgent 발표 PPT 생성
python tools/make_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 ─────────────────────────────────────────────────────────────────────
C_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
C_DARK    = RGBColor(0x33, 0x33, 0x33)
C_GRAY    = RGBColor(0x66, 0x66, 0x66)
C_LGRAY   = RGBColor(0xAA, 0xAA, 0xAA)
C_RULE    = RGBColor(0xDD, 0xDD, 0xDD)
C_BOX     = RGBColor(0xF5, 0xF5, 0xF5)
C_BOX2    = RGBColor(0xEB, 0xEB, 0xEB)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT  = RGBColor(0x1A, 0x1A, 0x1A)   # 강조 = 검정

FONT = "Pretendard"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── 유틸 ─────────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=C_BOX, line_color=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h,
        size=14, bold=False, color=C_DARK,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name   = FONT
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb


def hline(slide, y, x=Inches(0.8), w=None, color=C_RULE, thickness=0.75):
    if w is None:
        w = W - Inches(1.6)
    s = slide.shapes.add_shape(1, x, y, w, Pt(thickness))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def slide_header(slide, num, topic, time_label):
    """슬라이드 상단 메타 라인"""
    txt(slide, f"{num:02d} / 05",
        Inches(0.8), Inches(0.3), Inches(2.0), Inches(0.35),
        size=11, color=C_LGRAY)
    txt(slide, topic,
        Inches(0.8), Inches(0.3), Inches(10.0), Inches(0.35),
        size=11, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)
    txt(slide, time_label,
        W - Inches(2.4), Inches(0.3), Inches(1.6), Inches(0.35),
        size=11, color=C_LGRAY, align=PP_ALIGN.RIGHT)
    hline(slide, Inches(0.72))


def section_label(slide, text, y):
    txt(slide, text.upper(),
        Inches(0.8), y, Inches(4.0), Inches(0.28),
        size=9, bold=True, color=C_LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# 01 — 소개
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
rect(s1, 0, 0, W, H, fill=C_WHITE)
slide_header(s1, 1, "소개", "20초")

txt(s1, "DQAgent",
    Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.6),
    size=72, bold=True, color=C_BLACK, align=PP_ALIGN.LEFT)

txt(s1, "Data Quality Multi-Agent Pipeline",
    Inches(0.8), Inches(2.9), Inches(11.5), Inches(0.55),
    size=22, color=C_GRAY)

hline(s1, Inches(3.6), color=C_RULE)

txt(s1,
    "DuckDB 규칙 검증 + 결정론적 정규화 + Claude CLI 이상치 판단을 "
    "단일 LCEL 파이프라인으로 연결해 데이터 품질 보고서를 자동 생성합니다.",
    Inches(0.8), Inches(3.75), Inches(10.0), Inches(0.9),
    size=16, color=C_DARK)

txt(s1, "LangChain LCEL   /   DuckDB   /   Claude CLI   /   Python",
    Inches(0.8), Inches(4.75), Inches(11.5), Inches(0.38),
    size=12, color=C_LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# 02 — 설계 의도
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
rect(s2, 0, 0, W, H, fill=C_WHITE)
slide_header(s2, 2, "설계 의도", "1분")

txt(s2, "설계 의도",
    Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.65),
    size=32, bold=True, color=C_BLACK)

# 세 구역
cols = [
    ("어떤 문제를 풀려고 했는가", [
        "수동 검토는 반복적이고 기준이 사람마다 다르다",
        "LLM에 전체 레코드를 넘기면 비용이 선형으로 증가한다",
        "단일 규칙 엔진은 통계적 이상치를 판단할 수 없다",
        "변경 이유가 기록되지 않아 감사 추적이 불가능하다",
    ]),
    ("대상 사용자", [
        "정기적으로 고객·운영 데이터를 검토하는 데이터 분석가",
        "DQ 검사를 파이프라인에 통합하려는 데이터 엔지니어",
        "변경 이력을 보고서로 제출해야 하는 QA·감사 담당자",
        "LLM 에이전트 구조를 연구하는 학생·연구자",
    ]),
    ("기대 효과", [
        "50건 기준 수동 60분 -> 자동 95초",
        "LLM 호출을 이상치 레코드에만 한정해 비용 절감",
        "Changelog로 필드·레코드·이유 완전 추적",
        "3개 소스 합의로 단일 판단 대비 신뢰도 향상",
    ]),
]

col_w = Inches(3.7)
col_x = [Inches(0.8), Inches(4.75), Inches(8.7)]

for i, (heading, items) in enumerate(cols):
    cx = col_x[i]
    txt(s2, heading,
        cx, Inches(1.85), col_w, Inches(0.45),
        size=13, bold=True, color=C_BLACK)
    hline(s2, Inches(2.35), x=cx, w=col_w, color=C_RULE)
    for j, item in enumerate(items):
        txt(s2, f"- {item}",
            cx, Inches(2.5) + Inches(0.72) * j, col_w, Inches(0.65),
            size=12, color=C_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# 03 — Harness 구성
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
rect(s3, 0, 0, W, H, fill=C_WHITE)
slide_header(s3, 3, "Harness 구성", "1분")

txt(s3, "Harness 구성",
    Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.65),
    size=32, bold=True, color=C_BLACK)

components = [
    ("CLAUDE.md",
     "프로젝트 컨텍스트 문서",
     [
         "Claude Code가 세션 시작 시 자동으로 읽는 프로젝트 설명서",
         "ANTHROPIC_API_KEY 제외 이유, 2A/2B 분리 원칙 등 핵심 제약 명시",
         "새 규칙 추가 위치(Stage 1 R9+), 테스트 방법, 산출물 경로 포함",
     ]),
    ("Skills",
     "에이전트 역할 정의 파일",
     [
         "skills/ 디렉토리에 스테이지별 마크다운으로 역할·입출력 명세",
         "01_input -> 02_preprocessor -> 03_dq_llm -> 04_dq_code -> 05_report",
         "Claude가 각 에이전트의 책임 범위를 명확히 구분하도록 안내",
     ]),
    ("Hooks",
     "이벤트 기반 자동 실행 명령",
     [
         "SessionStart: 환경 점검 (Claude CLI, .env, report/ 디렉토리)",
         "PostToolUse Write|Edit: .py 파일 저장 시 py_compile 문법 검사",
         "PreToolUse Bash: 명령에 ANTHROPIC_API_KEY 포함 시 확인 요청",
     ]),
    ("MCP",
     "외부 도구 연결 인터페이스",
     [
         "현재 미구성 (이번 프로젝트 범위 외)",
         "확장 방향: DuckDB MCP로 SQL 직접 실행, 외부 DB 연결",
         "데이터 소스 다양화 시 MCP 서버 추가 예정",
     ]),
]

row_h = Inches(1.42)
row_y = Inches(1.75)
total_w = W - Inches(1.6)

for i, (label, sublabel, items) in enumerate(components):
    ry = row_y + row_h * i
    # 왼쪽 라벨 영역
    rect(s3, Inches(0.8), ry, Inches(2.4), row_h - Inches(0.06),
         fill=C_BOX if i % 2 == 0 else C_BOX2)
    txt(s3, label,
        Inches(0.88), ry + Inches(0.18), Inches(2.2), Inches(0.38),
        size=14, bold=True, color=C_BLACK)
    txt(s3, sublabel,
        Inches(0.88), ry + Inches(0.55), Inches(2.2), Inches(0.35),
        size=10, color=C_GRAY)
    # 오른쪽 내용
    for j, item in enumerate(items):
        txt(s3, f"- {item}",
            Inches(3.45), ry + Inches(0.15) + Inches(0.38) * j,
            Inches(9.6), Inches(0.36),
            size=11, color=C_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# 04 — 시연
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
rect(s4, 0, 0, W, H, fill=C_WHITE)
slide_header(s4, 4, "시연", "2분 20초")

txt(s4, "시연",
    Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.65),
    size=32, bold=True, color=C_BLACK)

# Demo 1
section_label(s4, "Demo 1  —  전체 파이프라인 실행", Inches(1.78))

rect(s4, Inches(0.8), Inches(2.08), Inches(11.5), Inches(0.72),
     fill=RGBColor(0x1A, 0x1A, 0x1A))
txt(s4, "python main.py customer_data_quality_test_50.json --batch-size 100 --skip-openai",
    Inches(1.0), Inches(2.13), Inches(11.0), Inches(0.6),
    size=12, color=C_WHITE, bold=False)

points_1 = [
    "Stage 1: DuckDB R1~R8 규칙 검증   ->   6건 위반 탐지 (Critical 4, Info 2)",
    "Stage 2A: 결정론적 정규화   ->   58건 처리 (LLM 호출 0회)",
    "Stage 2B: Claude 이상치 판단   ->   2건만 처리 (ambiguous_indices 라우팅)",
    "Stage 4: 최종 DQ 점수 94 / 100   A등급   /   소요 95초",
]
for j, p in enumerate(points_1):
    txt(s4, f"- {p}",
        Inches(0.8), Inches(2.92) + Inches(0.38) * j, Inches(11.5), Inches(0.35),
        size=11, color=C_DARK)

hline(s4, Inches(4.55), color=C_RULE)

# Demo 2
section_label(s4, "Demo 2  —  ROI 비교 분석", Inches(4.65))

rect(s4, Inches(0.8), Inches(4.95), Inches(11.5), Inches(0.72),
     fill=RGBColor(0x1A, 0x1A, 0x1A))
txt(s4,
    "python tools/roi_pandas.py customer_data_quality_test_50.json report/..._report.json --export roi_result.xlsx",
    Inches(1.0), Inches(5.0), Inches(11.0), Inches(0.6),
    size=12, color=C_WHITE)

points_2 = [
    "필드별 NULL 변화, 수치 통계 before/after, 액션 분포, Stage 2A vs 2B 기여도",
    "수동 60분 대비 Agent 95초   ->   ROI 29,900%   /   Excel 6-sheet 자동 저장",
]
for j, p in enumerate(points_2):
    txt(s4, f"- {p}",
        Inches(0.8), Inches(5.82) + Inches(0.4) * j, Inches(11.5), Inches(0.36),
        size=11, color=C_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# 05 — 회고
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
rect(s5, 0, 0, W, H, fill=C_WHITE)
slide_header(s5, 5, "회고", "20초")

txt(s5, "회고",
    Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.65),
    size=32, bold=True, color=C_BLACK)

retro = [
    ("잘된 점", C_BLACK, [
        "Stage 2A / 2B 분리로 LLM 호출을 50건에서 2건으로 줄였다",
        "Changelog 덕분에 무엇을 왜 바꿨는지 완전히 추적된다",
        "DuckDB 기반 프로파일링이 빠르고 SQL로 확장하기 쉽다",
        "파이프라인 실행이 명령어 한 줄로 끝난다",
    ]),
    ("어려웠던 점", C_BLACK, [
        "Git Bash subprocess에서 한국어 프롬프트 인코딩 문제 (cp949 오류)",
        "ANTHROPIC_API_KEY가 .env에 있으면 OAuth 대신 크레딧을 소모",
        "Claude CLI 타임아웃 -> 병렬 워커 수 조정으로 해결했지만 불안정",
        "pandas float NaN이 None으로 인식되지 않아 is_null() 헬퍼 필요",
    ]),
    ("개선하고 싶은 점", C_BLACK, [
        "대용량 처리: 워커 수를 데이터 크기에 따라 동적으로 조정",
        "실시간 처리: 현재 배치 전용 -> 스트리밍 지원 추가",
        "도메인 규칙: 의료·금융 등 업종별 특화 규칙 플러그인 구조",
        "MCP 연동: DuckDB MCP 서버로 외부 DB 직접 연결",
    ]),
]

col_w = Inches(3.7)
col_x = [Inches(0.8), Inches(4.75), Inches(8.7)]

for i, (heading, _, items) in enumerate(retro):
    cx = col_x[i]
    txt(s5, heading,
        cx, Inches(1.85), col_w, Inches(0.45),
        size=13, bold=True, color=C_BLACK)
    hline(s5, Inches(2.35), x=cx, w=col_w, color=C_RULE)
    for j, item in enumerate(items):
        txt(s5, f"- {item}",
            cx, Inches(2.5) + Inches(0.72) * j, col_w, Inches(0.65),
            size=11, color=C_DARK)


# ── 저장 ─────────────────────────────────────────────────────────────────────
out = "DQAgent_발표.pptx"
prs.save(out)
print(f"저장 완료: {out}")
